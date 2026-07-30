"""استخراج Document من ملف PowerPoint حديث (.pptx فقط، Office Open XML) عبر
python-pptx — مكتبة Python خالصة، بلا أي اعتمادية نظام (لا LibreOffice/سudo،
بنفس قيد البيئة الموثَّق في `ingest.py`). **لا يدعم .ppt القديم عمداً** (الصيغة
الثنائية القديمة تحتاج محرّك تحويل خارجي غير متاح هنا).

**قيد معماري مهم:** خلاف PDF، لا توجد وسيلة لعرض الشريحة الأصلية بصرياً في لوحة
الـPDF بالواجهة (python-pptx يقرأ بنية الشرائح فقط، لا يُصيّرها كصورة) — الواجهة
تعرض المحرر فقط بلا لوحة معاينة مطابقة للمصدر الأصلي لملفات PowerPoint (انظر
`page.tsx`، شرط `file.name.endsWith(".pdf")` قبل تركيب `PdfPane`)."""

from __future__ import annotations

import base64
from io import BytesIO
from typing import List, Optional

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .ingest import _sort_blocks_by_position
from .schema import Block, BlockType, BoundingBox, Document, ImageAsset, Page, PageSource, SourceEngine


def _shape_bbox(shape) -> Optional[BoundingBox]:
    """إحداثيات الشكل بوحدة EMU الخام (لا تحويل لنقاط/بكسل) — تكفي لترتيب القراءة
    النسبي (`_sort_blocks_by_position` يقارن مواضع فقط، لا يعتمد وحدة قياس مطلقة)،
    ولا معنى لتحويلها لأي وحدة أخرى بما أنه لا توجد لوحة PDF لمزامنتها معها أصلاً."""
    try:
        return BoundingBox(x0=shape.left, y0=shape.top, x1=shape.left + shape.width, y1=shape.top + shape.height)
    except (TypeError, AttributeError):
        return None


def _table_rows(shape) -> List[List[str]]:
    return [[cell.text for cell in row.cells] for row in shape.table.rows]


def _slide_blocks_and_images(slide, slide_number: int, images: List[ImageAsset]) -> List[Block]:
    blocks: List[Block] = []
    for shape in slide.shapes:
        bbox = _shape_bbox(shape)

        if shape.has_table:
            rows = _table_rows(shape)
            blocks.append(
                Block(block_type=BlockType.TABLE, rows=rows, raw_rows=rows, bbox=bbox, source_engine=SourceEngine.PYMUPDF)
            )
            continue

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                png_buffer = BytesIO()
                image = Image.open(BytesIO(shape.image.blob)).convert("RGB")
                image.save(png_buffer, format="PNG")
            except Exception:
                continue  # صورة تالفة/ترميز غير مدعوم — تفشل صورة واحدة، لا الشريحة كلها

            image_id = f"Image_{len(images) + 1:02d}"
            images.append(
                ImageAsset(
                    page_number=slide_number,
                    index=len(images),
                    image_id=image_id,
                    mime_type="image/png",
                    data_base64=base64.b64encode(png_buffer.getvalue()).decode("ascii"),
                    width=image.width,
                    height=image.height,
                    bbox=bbox,
                )
            )
            blocks.append(
                Block(
                    block_type=BlockType.PARAGRAPH,
                    text=f"[Insert {image_id} here]",
                    bbox=bbox,
                    source_engine=SourceEngine.PYMUPDF,
                )
            )
            continue

        if shape.has_text_frame:
            text = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text).strip()
            if not text:
                continue
            blocks.append(Block(block_type=BlockType.PARAGRAPH, text=text, bbox=bbox, source_engine=SourceEngine.PYMUPDF))

    return _sort_blocks_by_position(blocks)


def extract_pptx_document(pptx_path: str, file_name: Optional[str] = None) -> Document:
    """كل شريحة → Page واحدة (source=DIGITAL، لا OCR إطلاقاً — بنية مستخرَجة مباشرة
    من ملف XML). نص/جداول/صور تُستخرَج بنفس مبدأ الصفحة الرقمية في `ingest.py`
    (صورة Picture → ImageAsset + placeholder نصي في مكانها، بدل تضمينها Base64
    داخل الفقرة نفسها)."""
    presentation = Presentation(pptx_path)
    pages: List[Page] = []
    images: List[ImageAsset] = []

    for index, slide in enumerate(presentation.slides):
        blocks = _slide_blocks_and_images(slide, index + 1, images)
        pages.append(Page(page_number=index + 1, source=PageSource.DIGITAL, blocks=blocks))

    return Document(file_name=file_name or pptx_path, pages=pages, images=images)
