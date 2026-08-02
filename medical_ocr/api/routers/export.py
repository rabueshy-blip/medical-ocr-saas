"""تصدير محتوى المحرر (JSON من TipTap، بعد تحرير المترجم) إلى Word/PowerPoint حقيقيَّين
أو PDF — جداول كجداول حقيقية وليست نصاً مسطَّحاً. الصور تُضمَّن كبيانات حقيقية في مكانها
الأصلي داخل الملف نفسه (Word/PowerPoint/PDF عبر `attrs.src`، انظر `_add_image`) — قرار
مُحدَّث (كان placeholder نصي فقط + مجلد `images/` منفصل، طُلب لاحقاً تضمين الصورة
الفعلية مباشرة). الصورة تُسلَّم أيضاً في مجلد `images/` ضمن ZIP عند وجود صور
(`_stream_with_optional_images`، تُستخدَم من `export_docx` و`export_pptx` معاً) كنسخة
إضافية بدقة كاملة لمن يحتاجها منفصلة.

يقبل المحتوى المُحرَّر من الواجهة مباشرة (وليس Document الأصلي من extract-document) عمداً:
الهدف تصدير النتيجة *بعد* ترجمة/تعديل المترجم، لا النص الخام المُستخرَج.

تصدير PDF عبر Playwright (Chromium headless) وليس WeasyPrint عمداً — WeasyPrint يحتاج
Pango/Cairo على مستوى النظام غير المتاحَين في هذه البيئة (لا Homebrew/sudo)، بينما ثنائي
Chromium الخاص بـ Playwright يُنزَّل عبر pip بلا صلاحيات إدارية."""

from __future__ import annotations

import base64
import html
import io
import logging
import re
import zipfile
from typing import List, Optional

from docx import Document as DocxDocument
from docx.shared import Inches as DocxInches
from fastapi import APIRouter, HTTPException, Request
from fastapi.responses import StreamingResponse
from PIL import Image as PILImage
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches
from pydantic import BaseModel, Field

from ..rate_limit import limiter

logger = logging.getLogger(__name__)

router = APIRouter(tags=["export"])


class ExportImage(BaseModel):
    """صورة مُستخرَجة أصلاً عبر `extract-document` (`schema.ImageAsset`)، تُرسَل هنا فقط
    عند التصدير كي تُحزَم أيضاً كنسخة منفصلة في مجلد `images/` داخل ملف ZIP النهائي —
    الصورة نفسها مُضمَّنة أصلاً في موضعها داخل الملف الأساسي (Word/PDF/PPT) عبر
    `attrs.src` لعقدة الصورة في محتوى TipTap، انظر `_add_image`."""

    image_id: str
    mime_type: str = "image/png"
    data_base64: str


class ExportRequest(BaseModel):
    content: dict
    file_name: str = "translated_document"
    images: List[ExportImage] = Field(default_factory=list)


def _safe_file_name(file_name: str) -> str:
    """يبقي فقط أحرف ASCII آمنة — `\\w` في Python يطابق أحرفاً عربية/يونيكود أيضاً
    (ليس ASCII فقط)، فيمر اسم ملف عربي دون تغيير ثم يُسبِّب `UnicodeEncodeError` عند
    وضعه في ترويسة `Content-Disposition` (يجب أن تكون latin-1 قابلة للترميز حسب
    HTTP). اسم الملف الأصلي (بأي لغة) يبقى محفوظاً في المستند نفسه، هذا فقط لاسم
    ملف التنزيل."""
    return re.sub(r"[^a-zA-Z0-9_\-. ]", "_", file_name).strip() or "translated_document"


def _extract_text(node: dict) -> str:
    if node.get("type") == "text":
        return node.get("text", "")
    return "".join(_extract_text(child) for child in node.get("content", []))


def _add_paragraph_or_heading(doc: DocxDocument, node: dict) -> None:
    text = _extract_text(node)
    if node["type"] == "heading":
        level = node.get("attrs", {}).get("level") or 1
        doc.add_heading(text, level=min(max(int(level), 1), 9))
    else:
        doc.add_paragraph(text)


def _add_table(doc: DocxDocument, node: dict) -> None:
    """يبني جدول Word حقيقياً — بما فيها دمج الخلايا (colspan) القادمة من TipTap
    (`attrs.colspan` على tableCell/tableHeader، مُولَّدة أصلاً من خلايا مدمجة حقيقية في
    PDF المصدر عبر documentToTiptap.ts، أو من دمج يدوي للمترجم داخل المحرر).

    عرض الشبكة الكلي (`num_cols`) يُحسَب من **مجموع** colspan لكل صف (وليس عدد عناصر
    المحتوى) لأن صفاً فيه خلية مدمجة يحتوي عناصر JSON أقل من عرض الشبكة الفعلي."""
    row_nodes = node.get("content", [])
    if not row_nodes:
        return
    num_cols = max(
        (
            sum(
                int(cell_node.get("attrs", {}).get("colspan", 1))
                for cell_node in row_node.get("content", [])
            )
            for row_node in row_nodes
        ),
        default=0,
    )
    if num_cols == 0:
        return
    table = doc.add_table(rows=0, cols=num_cols)
    table.style = "Table Grid"
    for row_node in row_nodes:
        row_cells = table.add_row().cells
        col_cursor = 0
        for cell_node in row_node.get("content", []):
            if col_cursor >= num_cols:
                break
            colspan = int(cell_node.get("attrs", {}).get("colspan", 1))
            colspan = max(1, min(colspan, num_cols - col_cursor))
            target_cell = row_cells[col_cursor]
            target_cell.text = _extract_text(cell_node)
            if colspan > 1:
                target_cell.merge(row_cells[col_cursor + colspan - 1])
            col_cursor += colspan


_DATA_URL_PATTERN = re.compile(r"^data:[^;]+;base64,(?P<data>.+)$", re.DOTALL)

# عرض معقول لصورة مُضمَّنة في Word/PowerPoint — أكبر من هذا يفيض عادة عن حواف الصفحة/
# الشريحة، أصغر منه غير مقروء لصورة سريرية (أشعة/جدول ممسوح). الارتفاع يُحسَب تلقائياً
# بنفس نسبة الأبعاد الأصلية عبر add_picture(width=...) وحدها.
_EMBEDDED_IMAGE_MAX_WIDTH_INCHES = 5.5


def _decode_data_url(src: str) -> Optional[bytes]:
    """يفكّ base64 من data URL (نفس صيغة `imageAssetSrc` في الواجهة:
    `data:{mime};base64,{data}`). يُرجع None لأي src غير بصيغة data URL متوقَّعة
    (رابط خارجي مثلاً) بدل رفع استثناء — العنصر يتدهور إلى placeholder نصي في هذه الحالة."""
    match = _DATA_URL_PATTERN.match(src or "")
    if not match:
        return None
    try:
        return base64.b64decode(match.group("data"))
    except (base64.binascii.Error, ValueError):  # type: ignore[attr-defined]
        return None


def _picture_width_inches(image_bytes: bytes) -> float:
    """يحسب عرض الإدراج بالإنش: بافتراض 150 DPI (دقة نموذجية لصفحاتنا الممسوحة)
    مع سقف `_EMBEDDED_IMAGE_MAX_WIDTH_INCHES` — **لا يكبِّر** صورة صغيرة (شعار/أيقونة)
    فوق حجمها الطبيعي، فقط يمنع صورة ضخمة من الفيضان عن حواف الصفحة/الشريحة."""
    try:
        with PILImage.open(io.BytesIO(image_bytes)) as pil_image:
            pixel_width = pil_image.width
    except Exception:
        return _EMBEDDED_IMAGE_MAX_WIDTH_INCHES
    return min(pixel_width / 150, _EMBEDDED_IMAGE_MAX_WIDTH_INCHES)


def _add_image(doc: DocxDocument, node: dict) -> None:
    """يُضمِّن الصورة الحقيقية في مكانها (طلب مُحدَّث: العميل يريد ملف Word/PDF/PPT جاهزاً
    بصرياً بالصور في موضعها الأصلي، وليس فقط نصاً للمترجم) — `attrs.src` هو data URL
    (نفس تنسيق `imageAssetSrc` في الواجهة)، سواء أتى من الـplaceholder التلقائي
    (`documentToTiptap.ts` يحوّله الآن لعقدة صورة حقيقية) أو من سحب يدوي من مكتبة
    الوسائط. **تدهور آمن:** src مفقود/غير قابل للفك (حالة غير متوقَّعة) يكتب نفس
    عبارة الـplaceholder النصي القديمة بدل فشل التصدير بالكامل."""
    attrs = node.get("attrs", {})
    image_id = attrs.get("imageId")
    image_bytes = _decode_data_url(attrs.get("src", ""))
    if image_bytes is None:
        doc.add_paragraph(f"[Insert {image_id or 'Image'} here]")
        return
    doc.add_picture(io.BytesIO(image_bytes), width=DocxInches(_picture_width_inches(image_bytes)))


def _add_pptx_table(prs: Presentation, node: dict) -> None:
    """جدول PowerPoint حقيقي على شريحة مستقلة خاصة به (نفس منطق دمج colspan لـ`_add_table`
    لكن عبر `cell.merge()` الخاصة بجداول python-pptx بدل python-docx) — جدول وفقرات نص
    ممتزجان في نفس الشريحة يصيران غير مقروءَين بصرياً بسرعة، فكل جدول يأخذ شريحة فارغة
    (`slide_layouts[6]`) بلا عنوان يتنافس معه على المساحة."""
    row_nodes = node.get("content", [])
    if not row_nodes:
        return
    num_cols = max(
        (
            sum(int(cell_node.get("attrs", {}).get("colspan", 1)) for cell_node in row_node.get("content", []))
            for row_node in row_nodes
        ),
        default=0,
    )
    if num_cols == 0:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left, top = Inches(0.4), Inches(0.4)
    width, height = prs.slide_width - Inches(0.8), prs.slide_height - Inches(0.8)
    table = slide.shapes.add_table(len(row_nodes), num_cols, left, top, width, height).table

    for row_index, row_node in enumerate(row_nodes):
        col_cursor = 0
        for cell_node in row_node.get("content", []):
            if col_cursor >= num_cols:
                break
            colspan = int(cell_node.get("attrs", {}).get("colspan", 1))
            colspan = max(1, min(colspan, num_cols - col_cursor))
            table.cell(row_index, col_cursor).text = _extract_text(cell_node)
            if colspan > 1:
                table.cell(row_index, col_cursor).merge(table.cell(row_index, col_cursor + colspan - 1))
            col_cursor += colspan


def _add_pptx_image(prs: Presentation, image_bytes: bytes) -> None:
    """صورة على شريحة مستقلة خاصة بها (نفس نمط `_add_pptx_table`) — أبسط وأسلم من
    محاولة دسّها داخل `placeholders[1]` النصي لتخطيط "Title and Content"، ويضمن
    محاذاتها ومقاسها داخل حدود الشريحة مهما كانت أبعادها الأصلية."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    max_width = prs.slide_width - Inches(0.8)
    max_height = prs.slide_height - Inches(0.8)
    try:
        with PILImage.open(io.BytesIO(image_bytes)) as pil_image:
            aspect = pil_image.height / pil_image.width if pil_image.width else 1.0
    except Exception:
        aspect = 1.0
    width, height = max_width, int(max_width * aspect)
    if height > max_height:
        width, height = int(max_height / aspect) if aspect else max_width, max_height
    left, top = (prs.slide_width - width) // 2, (prs.slide_height - height) // 2
    slide.shapes.add_picture(io.BytesIO(image_bytes), left, top, width=width, height=height)


def _build_pptx(top_content: List[dict]) -> bytes:
    """يحوّل نفس محتوى TipTap المستخدَم في `export_docx` إلى عرض PowerPoint: كل
    عنوان (heading) يبدأ شريحة جديدة (تخطيط "Title and Content")، والفقرات
    التالية تتراكم كنقاط في مساحة المحتوى لنفس الشريحة حتى العنوان التالي أو
    الجدول التالي. الجدول يأخذ شريحته المستقلة دوماً (انظر `_add_pptx_table`)،
    وكذلك الصورة (انظر `_add_pptx_image`) — تُضمَّن حقيقة إن وُجد `attrs.src`
    قابل للفك (نفس منطق `_add_image` في Word)، وإلا نص placeholder كتدهور آمن.

    **قيد معروف مقبول:** خلافاً لـWord/PDF، الشريحة لا "تفيض" تلقائياً لشريحة
    تالية عند نص طويل جداً — تجاوز عملي غير محلول هنا، الشرائح الفعلية لهذا
    المشروع (تقارير طبية قصيرة الفقرات) لم تُظهر هذه المشكلة عملياً."""
    prs = Presentation()
    current_body = None

    def start_slide(title: str) -> None:
        nonlocal current_body
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        current_body = slide.placeholders[1].text_frame
        current_body.clear()

    def append_text(text: str) -> None:
        nonlocal current_body
        if current_body is None:
            start_slide("")
        if len(current_body.paragraphs) == 1 and not current_body.paragraphs[0].text:
            current_body.paragraphs[0].text = text
        else:
            current_body.add_paragraph().text = text

    for node in top_content:
        node_type = node.get("type")
        try:
            if node_type == "heading":
                start_slide(_extract_text(node))
            elif node_type == "paragraph":
                text = _extract_text(node)
                if text:
                    append_text(text)
            elif node_type == "table":
                _add_pptx_table(prs, node)
                current_body = None  # الجدول شريحة مستقلة؛ الفقرة التالية تبدأ شريحة نص جديدة
            elif node_type == "image":
                attrs = node.get("attrs", {})
                image_id = attrs.get("imageId")
                image_bytes = _decode_data_url(attrs.get("src", ""))
                if image_bytes is None:
                    append_text(f"[Insert {image_id or 'Image'} here]")
                else:
                    _add_pptx_image(prs, image_bytes)
                    current_body = None  # نفس منطق الجدول: الشريحة النصية التالية تبدأ من جديد
        except Exception as exc:  # عنصر واحد فاشل لا يوقف تصدير بقية العرض
            logger.warning("تعذّر تصدير عنصر PowerPoint من نوع %s: %s", node_type, exc)

    if len(prs.slides) == 0:
        start_slide("")

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def _stream_with_optional_images(
    primary_bytes: bytes,
    primary_extension: str,
    primary_media_type: str,
    safe_name: str,
    images: List[ExportImage],
) -> StreamingResponse:
    """يُرجع الملف الأساسي (.docx/.pptx) مباشرة إن لم توجد صور، أو ZIP واحد (الملف +
    مجلد `images/`) عند وجودها — نفس منطق `export_docx` الأصلي، مُستخرَج هنا كي
    يُشارَك مع `export_pptx` بدل تكراره."""
    if not images:
        return StreamingResponse(
            io.BytesIO(primary_bytes),
            media_type=primary_media_type,
            headers={"Content-Disposition": f'attachment; filename="{safe_name}.{primary_extension}"'},
        )

    # مستند فيه صور: يُسلَّم ZIP واحد (الملف + مجلد images/) بدل ملف مفرد، كي يجد
    # المترجم/فريق DTP الصور التي تشير إليها الـplaceholders النصية داخل الملف بسهولة.
    zip_buffer = io.BytesIO()
    with zipfile.ZipFile(zip_buffer, mode="w", compression=zipfile.ZIP_DEFLATED) as zip_file:
        zip_file.writestr(f"{safe_name}.{primary_extension}", primary_bytes)
        for image in images:
            extension = image.mime_type.rsplit("/", maxsplit=1)[-1] or "png"
            zip_file.writestr(f"images/{image.image_id}.{extension}", base64.b64decode(image.data_base64))
    zip_buffer.seek(0)

    return StreamingResponse(
        zip_buffer,
        media_type="application/zip",
        headers={"Content-Disposition": f'attachment; filename="{safe_name}.zip"'},
    )


@router.post("/export-docx")
@limiter.limit("20/minute")
def export_docx(request: Request, payload: ExportRequest) -> StreamingResponse:
    top_content = payload.content.get("content", [])
    if not top_content:
        raise HTTPException(status_code=422, detail="المستند فارغ، لا يوجد محتوى للتصدير")

    doc = DocxDocument()
    for node in top_content:
        node_type = node.get("type")
        try:
            if node_type in ("paragraph", "heading"):
                _add_paragraph_or_heading(doc, node)
            elif node_type == "table":
                _add_table(doc, node)
            elif node_type == "image":
                _add_image(doc, node)
        except Exception as exc:  # عنصر واحد فاشل لا يوقف تصدير بقية المستند
            logger.warning("تعذّر تصدير عنصر من نوع %s: %s", node_type, exc)

    docx_buffer = io.BytesIO()
    doc.save(docx_buffer)

    safe_name = _safe_file_name(payload.file_name)
    return _stream_with_optional_images(
        docx_buffer.getvalue(),
        "docx",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        safe_name,
        payload.images,
    )


@router.post("/export-pptx")
@limiter.limit("20/minute")
def export_pptx(request: Request, payload: ExportRequest) -> StreamingResponse:
    top_content = payload.content.get("content", [])
    if not top_content:
        raise HTTPException(status_code=422, detail="المستند فارغ، لا يوجد محتوى للتصدير")

    pptx_bytes = _build_pptx(top_content)

    safe_name = _safe_file_name(payload.file_name)
    return _stream_with_optional_images(
        pptx_bytes,
        "pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        safe_name,
        payload.images,
    )


# نفس سقف Word (150 DPI × 5.5 إنش) لكن هنا يُطبَّق فعلياً على أبعاد البكسل قبل
# التضمين، وليس فقط عرض عرض CSS — الفرق مهم: صورة مسحوبة بحجمها الأصلي (لوحظ حتى
# 1512×1378px لكل صورة، 9 صور في مستند واحد) تُجبر Chromium على فك وترسيم كل تلك
# البكسلات فعلياً رغم أن CSS `max-width:100%` يعرضها أصغر بصرياً فقط — على حاوية
# Render المجانية (512MB، وعملية الاستخراج نفسها لوحظت تصل لذروة RSS ~360MB لنفس
# الملف) هذا كان يُسقط/يُعلِّق Playwright فعلياً (طلب حقيقي بقي معلَّقاً على الإنتاج
# لأكثر من دقيقتين بلا استجابة على ملف حقيقي 8 صفحات/9 صور، بينما نفس التصدير محلياً
# استغرق أقل من 15 ثانية) — تصغير أبعاد البكسل الفعلية قبل الترسيم هو الإصلاح الجذري.
_PDF_EMBEDDED_IMAGE_MAX_WIDTH_PX = 825


def _resize_image_data_url(src: str) -> str:
    """يُصغِّر صورة data URL إلى `_PDF_EMBEDDED_IMAGE_MAX_WIDTH_PX` كحد أقصى قبل
    تمريرها لـChromium عبر Playwright — يُرجع `src` كما هو دون تعديل لأي صورة أصغر
    من الحد أصلاً، أو لأي src غير قابل للفك (تدهور آمن، نفس نمط `_decode_data_url`)."""
    image_bytes = _decode_data_url(src)
    if image_bytes is None:
        return src
    try:
        with PILImage.open(io.BytesIO(image_bytes)) as pil_image:
            if pil_image.width <= _PDF_EMBEDDED_IMAGE_MAX_WIDTH_PX:
                return src
            ratio = _PDF_EMBEDDED_IMAGE_MAX_WIDTH_PX / pil_image.width
            resized = pil_image.resize(
                (_PDF_EMBEDDED_IMAGE_MAX_WIDTH_PX, max(1, round(pil_image.height * ratio))),
                PILImage.LANCZOS,
            )
            buffer = io.BytesIO()
            resized.save(buffer, format="PNG")
            return f"data:image/png;base64,{base64.b64encode(buffer.getvalue()).decode('ascii')}"
    except Exception:
        return src


def _node_to_html(node: dict) -> str:
    node_type = node.get("type")
    if node_type == "text":
        return html.escape(node.get("text", ""))
    if node_type == "paragraph":
        inner = "".join(_node_to_html(child) for child in node.get("content", []))
        return f"<p>{inner}</p>" if inner else "<p>&nbsp;</p>"
    if node_type == "heading":
        level = min(max(int(node.get("attrs", {}).get("level") or 1), 1), 6)
        inner = "".join(_node_to_html(child) for child in node.get("content", []))
        return f"<h{level}>{inner}</h{level}>"
    if node_type == "table":
        rows_html = []
        for row in node.get("content", []):
            cells_html = []
            for cell in row.get("content", []):
                tag = "th" if cell.get("type") == "tableHeader" else "td"
                colspan = int(cell.get("attrs", {}).get("colspan", 1))
                colspan_attr = f' colspan="{colspan}"' if colspan > 1 else ""
                inner = "".join(_node_to_html(child) for child in cell.get("content", []))
                cells_html.append(f"<{tag}{colspan_attr}>{inner}</{tag}>")
            rows_html.append(f"<tr>{''.join(cells_html)}</tr>")
        return f"<table>{''.join(rows_html)}</table>"
    if node_type == "image":
        src = node.get("attrs", {}).get("src", "")
        if not src:
            return ""
        return f'<img src="{html.escape(_resize_image_data_url(src))}" />'
    return ""


def _content_to_html_document(content: dict) -> str:
    body = "".join(_node_to_html(node) for node in content.get("content", []))
    return f"""<!doctype html>
<html>
<head>
<meta charset="utf-8" />
<style>
  body {{ font-family: -apple-system, "Segoe UI", Tahoma, Arial, sans-serif;
          padding: 20px; line-height: 1.6; color: #111; }}
  table {{ border-collapse: collapse; width: 100%; margin: 12px 0; }}
  td, th {{ border: 1px solid #333; padding: 6px 10px; text-align: start; }}
  th {{ background: #f2f2f2; }}
  img {{ max-width: 100%; margin: 12px 0; display: block; }}
  p:empty::before {{ content: "\\00a0"; }}
</style>
</head>
<body>{body}</body>
</html>"""


@router.post("/export-pdf")
@limiter.limit("10/minute")
def export_pdf(request: Request, payload: ExportRequest) -> StreamingResponse:
    top_content = payload.content.get("content", [])
    if not top_content:
        raise HTTPException(status_code=422, detail="المستند فارغ، لا يوجد محتوى للتصدير")

    html_document = _content_to_html_document(payload.content)

    try:
        with sync_playwright() as playwright:
            browser = playwright.chromium.launch()
            try:
                page = browser.new_page()
                page.set_content(html_document, wait_until="load")
                pdf_bytes = page.pdf(
                    format="A4",
                    print_background=True,
                    margin={"top": "20mm", "bottom": "20mm", "left": "15mm", "right": "15mm"},
                )
            finally:
                browser.close()
    except Exception as exc:
        logger.error("فشل توليد PDF عبر Playwright: %s", exc)
        raise HTTPException(status_code=500, detail=f"فشل توليد PDF: {exc}") from exc

    safe_name = _safe_file_name(payload.file_name)
    return StreamingResponse(
        io.BytesIO(pdf_bytes),
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{safe_name}.pdf"'},
    )
