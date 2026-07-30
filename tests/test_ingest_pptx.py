"""اختبار استخراج Document من ملفات PowerPoint حديثة (.pptx) — ميزة جديدة طلبها
المستخدم صراحة، مقصورة عمداً على .pptx (Office Open XML) بلا دعم .ppt القديم
(يحتاج LibreOffice/أداة تحويل نظام غير متاحة في بيئة المشروع، انظر توثيق
`medical_ocr/ingest_pptx.py`)."""

import os
import tempfile
import unittest
from io import BytesIO

from fastapi.testclient import TestClient
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from medical_ocr.api.app import app
from medical_ocr.ingest_pptx import extract_pptx_document
from medical_ocr.schema import BlockType


def _make_pptx(path: str) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # تخطيط فارغ تماماً

    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(1))
    textbox.text_frame.text = "Slide title text"

    table_shape = slide.shapes.add_table(2, 2, Inches(0.5), Inches(2), Inches(3), Inches(1))
    table_shape.table.cell(0, 0).text = "Header1"
    table_shape.table.cell(0, 1).text = "Header2"
    table_shape.table.cell(1, 0).text = "Value1"
    table_shape.table.cell(1, 1).text = "Value2"

    image = Image.new("RGB", (100, 80), (200, 50, 50))
    buffer = BytesIO()
    image.save(buffer, format="PNG")
    buffer.seek(0)
    slide.shapes.add_picture(buffer, Inches(0.5), Inches(4), Inches(2), Inches(1.5))

    prs.save(path)


class TestExtractPptxDocument(unittest.TestCase):
    def test_extracts_text_table_and_image_with_placeholder(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            pptx_path = os.path.join(tmp_dir, "slide.pptx")
            _make_pptx(pptx_path)

            document = extract_pptx_document(pptx_path, file_name="slide.pptx")

            self.assertEqual(len(document.pages), 1)
            page = document.pages[0]

            text_blocks = [b for b in page.blocks if b.block_type == BlockType.PARAGRAPH and b.text != "[Insert Image_01 here]"]
            table_blocks = [b for b in page.blocks if b.block_type == BlockType.TABLE]
            placeholder_blocks = [b for b in page.blocks if b.text == "[Insert Image_01 here]"]

            self.assertEqual([b.text for b in text_blocks], ["Slide title text"])
            self.assertEqual(len(table_blocks), 1)
            self.assertEqual(table_blocks[0].rows, [["Header1", "Header2"], ["Value1", "Value2"]])
            self.assertEqual(len(placeholder_blocks), 1)

            self.assertEqual(len(document.images), 1)
            asset = document.images[0]
            self.assertEqual(asset.image_id, "Image_01")
            self.assertEqual(asset.mime_type, "image/png")
            self.assertEqual((asset.width, asset.height), (100, 80))

    def test_empty_presentation_yields_document_with_no_blocks(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            pptx_path = os.path.join(tmp_dir, "empty.pptx")
            prs = Presentation()
            prs.slides.add_slide(prs.slide_layouts[6])
            prs.save(pptx_path)

            document = extract_pptx_document(pptx_path, file_name="empty.pptx")

            self.assertEqual(len(document.pages), 1)
            self.assertEqual(document.pages[0].blocks, [])
            self.assertEqual(document.images, [])


class TestExtractPptxEndpoint(unittest.TestCase):
    def setUp(self):
        self.client = TestClient(app)
        self.headers = {}
        api_key = os.environ.get("APP_API_KEY")
        if api_key:
            self.headers["X-API-Key"] = api_key

    def test_rejects_non_pptx_extension(self):
        response = self.client.post(
            "/extract-pptx",
            files={"file": ("not_a_pptx.txt", b"plain text", "text/plain")},
            headers=self.headers,
        )
        self.assertEqual(response.status_code, 422)

    def test_extracts_uploaded_pptx_via_http(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            pptx_path = os.path.join(tmp_dir, "slide.pptx")
            _make_pptx(pptx_path)

            with open(pptx_path, "rb") as f:
                response = self.client.post(
                    "/extract-pptx",
                    files={
                        "file": (
                            "slide.pptx",
                            f.read(),
                            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        )
                    },
                    headers=self.headers,
                )

        self.assertEqual(response.status_code, 200)
        body = response.json()
        self.assertEqual(len(body["images"]), 1)
        self.assertEqual(len(body["pages"]), 1)


if __name__ == "__main__":
    unittest.main()
